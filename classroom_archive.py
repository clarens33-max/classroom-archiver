"""
classroom_archive.py
────────────────────────────────────────────────────────────────────────────────
Archives a Google Classroom course using the Classroom, Drive, and Slides APIs.

For each lesson it:
  - Exports Google Slides decks as PNG images → stitched into a local .pptx
  - Extracts all hyperlinks from slides → saved to a _links.txt file
  - Downloads all attached Drive files (Docs → PDF, Sheets → xlsx, etc.)
  - Saves external URLs and chat logs
  - Parses Vimeo URLs AND passwords from material descriptions automatically
  - Writes a _vimeo_queue.json for classroom_transcribe.py to process

Output structure:
  output/
    <CourseName>/
      Lesson_01_<Title>/
        slides_<DeckName>.pptx
        slides_<DeckName>_links.txt
        resources/
          <file>.pdf
        resources_urls.txt
        description.txt        ← assignment brief if present
      _vimeo_queue.json        ← feed into classroom_transcribe.py

Prerequisites:
  pip install google-auth google-auth-oauthlib google-api-python-client
              python-pptx Pillow requests
  credentials.json in the same folder (see google_cloud_setup.md)
────────────────────────────────────────────────────────────────────────────────
"""

import os
import io
import re
import json
import time
from pathlib import Path

import requests
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from google.auth.transport.requests import Request
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload

# ── CONFIG ────────────────────────────────────────────────────────────────────

OUTPUT_DIR = "output"

SCOPES = [
    "https://www.googleapis.com/auth/classroom.courses.readonly",
    "https://www.googleapis.com/auth/classroom.coursework.me.readonly",
    "https://www.googleapis.com/auth/classroom.courseworkmaterials.readonly",
    "https://www.googleapis.com/auth/classroom.announcements.readonly",
    "https://www.googleapis.com/auth/drive",
]

# Matches vimeo.com and player.vimeo.com URLs
VIMEO_RE = re.compile(
    r"https?://(?:www\.|player\.)?vimeo\.com/(?:video/)?(\d+)[^\s\"'<]*"
)

# Extracts password from strings like "Password: AISA26/11" or "password: abc123"
PASSWORD_RE = re.compile(
    r"[Pp]assword[:\s]+([A-Za-z0-9@#$%^&*!_/\-\.]+)"
)

# ── AUTH ──────────────────────────────────────────────────────────────────────

def get_services():
    os.environ["OAUTHLIB_RELAX_TOKEN_SCOPE"] = "1"
    creds = None
    if os.path.exists("token.json"):
        creds = Credentials.from_authorized_user_file("token.json", SCOPES)
    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            creds.refresh(Request())
        else:
            flow = InstalledAppFlow.from_client_secrets_file(
                "credentials.json", SCOPES
            )
            creds = flow.run_local_server(port=0)
        with open("token.json", "w") as f:
            f.write(creds.to_json())
    classroom = build("classroom", "v1", credentials=creds)
    drive     = build("drive",     "v3", credentials=creds)
    slides    = build("slides",    "v1", credentials=creds)
    return classroom, drive, slides, creds


# ── HELPERS ───────────────────────────────────────────────────────────────────

def safe_name(s, max_len=60):
    s = re.sub(r'[\\/*?:"<>|]', "_", str(s))
    s = re.sub(r"\s+", "_", s.strip())
    return s[:max_len]


def extract_vimeo_info(text):
    """Return list of (url, password) tuples found in a block of text."""
    if not text:
        return []
    urls      = VIMEO_RE.findall(text)  # returns video IDs
    url_full  = VIMEO_RE.findall(text)
    passwords = PASSWORD_RE.findall(text)

    # Re-extract full URLs (findall on a group loses the full match)
    full_urls = [m.group(0) for m in VIMEO_RE.finditer(text)]
    password  = passwords[0] if passwords else ""

    return [(url, password) for url in full_urls]


# ── SLIDES EXPORT ─────────────────────────────────────────────────────────────

def export_slides(slides_svc, creds, file_id, deck_name, out_dir):
    """Export a Google Slides deck as PNG thumbnails → local .pptx"""
    print(f"    📊 Exporting slides: {deck_name}")
    try:
        pres = slides_svc.presentations().get(
            presentationId=file_id
        ).execute()
    except Exception as e:
        print(f"    ❌ Cannot access slides: {e}")
        return []

    slide_list = pres.get("slides", [])
    page_size  = pres.get("pageSize", {})
    w_in = page_size.get("width",  {}).get("magnitude", 9144000) / 914400
    h_in = page_size.get("height", {}).get("magnitude", 5143500) / 914400

    prs = Presentation()
    prs.slide_width  = Inches(w_in)
    prs.slide_height = Inches(h_in)
    blank = prs.slide_layouts[6]

    headers = {"Authorization": f"Bearer {creds.token}"}
    links   = []

    for idx, slide in enumerate(slide_list, 1):
        slide_id = slide["objectId"]

        # Thumbnail
        thumb_resp = requests.get(
            f"https://slides.googleapis.com/v1/presentations/{file_id}"
            f"/pages/{slide_id}/thumbnail?thumbnailProperties.thumbnailSize=LARGE",
            headers=headers,
        )
        if thumb_resp.status_code != 200:
            print(f"      ⚠️  Slide {idx}: thumbnail failed")
            continue

        img_url  = thumb_resp.json().get("contentUrl", "")
        img_data = io.BytesIO(requests.get(img_url).content)

        pptx_slide = prs.slides.add_slide(blank)
        pptx_slide.shapes.add_picture(
            img_data, Inches(0), Inches(0),
            width=Inches(w_in), height=Inches(h_in)
        )

        # Extract hyperlinks
        for el in slide.get("pageElements", []):
            _collect_links(el, idx, links)

        if idx % 5 == 0:
            print(f"      … {idx}/{len(slide_list)} slides")
        time.sleep(0.25)

    pptx_path = out_dir / f"{safe_name(deck_name)}.pptx"
    prs.save(str(pptx_path))
    print(f"    💾 Saved {len(slide_list)}-slide deck → {pptx_path.name}")

    if links:
        links_path = out_dir / f"{safe_name(deck_name)}_links.txt"
        with open(links_path, "w", encoding="utf-8") as f:
            f.write(f"Links extracted from: {deck_name}\n{'='*60}\n\n")
            for lnk in links:
                f.write(f"Slide {lnk['slide']:02d} | {lnk['text']:<40} → {lnk['url']}\n")
        print(f"    🔗 {len(links)} links → {links_path.name}")

    return links


def _collect_links(element, slide_idx, out):
    """Recursively collect hyperlinks from a slide page element."""
    for te in element.get("shape", {}).get("text", {}).get("textElements", []):
        url = te.get("textRun", {}).get("style", {}).get("link", {}).get("url", "")
        if url:
            text = te.get("textRun", {}).get("content", "").strip()[:40]
            out.append({"slide": slide_idx, "text": text, "url": url})
    for row in element.get("table", {}).get("tableRows", []):
        for cell in row.get("tableCells", []):
            for te in cell.get("text", {}).get("textElements", []):
                url = te.get("textRun", {}).get("style", {}).get("link", {}).get("url", "")
                if url:
                    out.append({"slide": slide_idx,
                                "text": te.get("textRun", {}).get("content", "")[:40].strip(),
                                "url": url})
    for child in element.get("elementGroup", {}).get("children", []):
        _collect_links(child, slide_idx, out)


# ── DRIVE FILE DOWNLOAD ────────────────────────────────────────────────────────

EXPORT_MAP = {
    "application/vnd.google-apps.document":
        ("application/pdf", ".pdf"),
    "application/vnd.google-apps.spreadsheet":
        ("application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", ".xlsx"),
    "application/vnd.google-apps.presentation":
        ("application/vnd.openxmlformats-officedocument.presentationml.presentation", ".pptx"),
    "application/vnd.google-apps.drawing":
        ("image/png", ".png"),
}

def download_drive_file(drive_svc, file_id, file_name, mime_type, out_dir):
    """Returns True on success, False if the file is inaccessible via API."""
    out_dir.mkdir(parents=True, exist_ok=True)
    if mime_type in EXPORT_MAP:
        export_mime, ext = EXPORT_MAP[mime_type]
        dest = out_dir / (safe_name(file_name) + ext)
        if dest.exists():
            print(f"      ⏭️  Exists: {dest.name}")
            return True
        try:
            req = drive_svc.files().export_media(
                fileId=file_id, mimeType=export_mime
            )
            buf = io.BytesIO()
            dl  = MediaIoBaseDownload(buf, req)
            done = False
            while not done:
                _, done = dl.next_chunk()
            dest.write_bytes(buf.getvalue())
            print(f"      💾 {dest.name}")
            return True
        except Exception:
            return False  # caller will handle the fallback
    else:
        dest = out_dir / safe_name(file_name)
        if dest.exists():
            print(f"      ⏭️  Exists: {dest.name}")
            return True
        try:
            req = drive_svc.files().get_media(fileId=file_id)
        except Exception as e:
            print(f"      ❌ Download failed for {file_name}: {e}")
            return False
        buf = io.BytesIO()
        dl  = MediaIoBaseDownload(buf, req)
        done = False
        try:
            while not done:
                _, done = dl.next_chunk()
        except Exception as e:
            print(f"      ❌ Download failed for {file_name}: {e}")
            return False
        dest.write_bytes(buf.getvalue())
        print(f"      💾 {dest.name}")
        return True


# ── ATTACHMENT PROCESSOR ──────────────────────────────────────────────────────

def process_attachments(attachments, lesson_dir, slides_svc, drive_svc,
                        creds, vimeo_queue, course_title, lesson_name,
                        extra_text=""):
    """Process a list of material attachments."""
    resources_dir = lesson_dir / "resources"
    external_urls = []

    # Also scan any free text (description body) for Vimeo links + passwords
    all_text = extra_text
    for att in attachments:
        if "link" in att:
            all_text += " " + att["link"].get("url", "")

    vimeo_hits = extract_vimeo_info(all_text)
    for url, password in vimeo_hits:
        vimeo_queue.append({
            "course":   course_title,
            "lesson":   lesson_name,
            "title":    lesson_name,
            "url":      url,
            "password": password,
        })
        print(f"      🎬 Vimeo queued: {url}  (pwd: {password or 'none'})")

    manual_downloads = []  # (file_name, url) for files inaccessible via API

    for att in attachments:
        if "driveFile" in att:
            # Normalise the nested structure
            df = att["driveFile"]
            if "driveFile" in df:
                df = df["driveFile"]
            file_id       = df.get("id", "")
            file_name     = df.get("title", file_id)
            alternate_url = df.get("alternateLink", f"https://drive.google.com/file/d/{file_id}")
            if not file_id:
                continue
            try:
                meta = drive_svc.files().get(
                    fileId=file_id, fields="mimeType,name"
                ).execute()
                mime      = meta.get("mimeType", "")
                file_name = meta.get("name", file_name)
            except Exception:
                mime = ""

            ok = True
            if mime == "application/vnd.google-apps.presentation":
                export_slides(slides_svc, creds, file_id, file_name, lesson_dir)
            elif mime == "":
                # Drive metadata inaccessible — try Slides API, then Doc PDF export
                print(f"      ⚠️  Can't read metadata for {file_name!r}, trying Slides API...")
                try:
                    pres = slides_svc.presentations().get(presentationId=file_id).execute()
                    export_slides(slides_svc, creds, file_id, file_name, lesson_dir)
                except Exception:
                    print(f"      ↩️  Not a Slides file, trying PDF export...")
                    ok = download_drive_file(drive_svc, file_id, file_name,
                                             "application/vnd.google-apps.document", resources_dir)
            else:
                ok = download_drive_file(drive_svc, file_id, file_name, mime, resources_dir)

            if not ok:
                print(f"      🔗 Inaccessible via API — saved to manual_download_needed.txt")
                manual_downloads.append((file_name, alternate_url))

        elif "link" in att:
            url   = att["link"].get("url", "")
            title = att["link"].get("title", url)
            # Vimeo already handled above via text scan — skip duplicates
            if VIMEO_RE.search(url):
                pass  # already queued above
            else:
                # Detect Google Workspace URLs and download via Drive API
                gdoc_match = re.search(
                    r"docs\.google\.com/(?:document|spreadsheets|presentation|drawings)/d/([a-zA-Z0-9_-]+)",
                    url,
                )
                if gdoc_match:
                    file_id = gdoc_match.group(1)
                    # Determine mime type from URL path
                    if "/document/" in url:
                        mime = "application/vnd.google-apps.document"
                    elif "/spreadsheets/" in url:
                        mime = "application/vnd.google-apps.spreadsheet"
                    elif "/presentation/" in url:
                        mime = "application/vnd.google-apps.presentation"
                    elif "/drawings/" in url:
                        mime = "application/vnd.google-apps.drawing"
                    else:
                        mime = "application/vnd.google-apps.document"
                    try:
                        meta = drive_svc.files().get(
                            fileId=file_id, fields="name"
                        ).execute()
                        file_name = meta.get("name", title)
                    except Exception:
                        file_name = title
                    print(f"      📄 Google Doc (link): {file_name}")
                    if mime == "application/vnd.google-apps.presentation":
                        export_slides(slides_svc, creds, file_id, file_name, lesson_dir)
                    else:
                        download_drive_file(drive_svc, file_id, file_name, mime, resources_dir)
                else:
                    external_urls.append(f"{title}  →  {url}")

        elif "youtubeVideo" in att:
            yv  = att["youtubeVideo"]
            external_urls.append(
                f"[YouTube] {yv.get('title','')}  →  https://youtu.be/{yv.get('id','')}"
            )

        elif "form" in att:
            fm  = att["form"]
            external_urls.append(
                f"[Form] {fm.get('title','')}  →  {fm.get('formUrl','')}"
            )

    if external_urls:
        resources_dir.mkdir(parents=True, exist_ok=True)
        urls_file = lesson_dir / "resources_urls.txt"
        with open(urls_file, "w", encoding="utf-8") as f:
            f.write("External resources\n" + "=" * 60 + "\n\n")
            f.write("\n".join(external_urls))
        print(f"      📋 {len(external_urls)} external URLs → resources_urls.txt")

    if manual_downloads:
        manual_file = lesson_dir / "manual_download_needed.txt"
        with open(manual_file, "w", encoding="utf-8") as f:
            f.write("Files inaccessible via API — open each URL in your browser to download or copy\n")
            f.write("=" * 70 + "\n\n")
            for fname, url in manual_downloads:
                f.write(f"{fname}\n  → {url}\n\n")
        print(f"      📎 {len(manual_downloads)} file(s) need manual download → manual_download_needed.txt")


# ── MAIN ──────────────────────────────────────────────────────────────────────

def main():
    print("🔐 Authenticating with Google...")
    classroom, drive, slides, creds = get_services()

    # Select course
    print("\n📚 Fetching your active courses...\n")
    resp    = classroom.courses().list(courseStates=["ACTIVE"]).execute()
    courses = resp.get("courses", [])
    if not courses:
        print("No active courses found.")
        return

    for i, c in enumerate(courses):
        print(f"  [{i+1}] {c['name']}")

    choice = input("\nEnter course number to archive: ").strip()
    try:
        course = courses[int(choice) - 1]
    except (ValueError, IndexError):
        print("Invalid selection.")
        return

    course_id    = course["id"]
    course_title = course["name"]
    print(f"\n✅ Archiving: {course_title}\n")

    base_dir = Path(OUTPUT_DIR) / safe_name(course_title)
    base_dir.mkdir(parents=True, exist_ok=True)

    # Load existing queue so re-runs don't lose previously found videos
    q_path = base_dir / "_vimeo_queue.json"
    existing_queue = json.loads(q_path.read_text(encoding="utf-8")) if q_path.exists() else []
    existing_urls  = {e["url"] for e in existing_queue}
    vimeo_queue = []

    # Fetch topics (used as lesson labels)
    topics = {}
    try:
        for t in classroom.courses().topics().list(
            courseId=course_id
        ).execute().get("topic", []):
            topics[t["topicId"]] = t["name"]
    except Exception:
        pass

    # ── Course materials ──────────────────────────────────────────────────────
    print("📂 Processing course materials...\n")
    try:
        mats = classroom.courses().courseWorkMaterials().list(
            courseId=course_id
        ).execute().get("courseWorkMaterial", [])
    except Exception as e:
        print(f"  ⚠️  Could not fetch materials: {e}")
        mats = []

    for idx, mat in enumerate(mats, 1):
        title      = mat.get("title", f"Material_{idx}")
        topic      = topics.get(mat.get("topicId", ""), "")
        prefix     = f"M{idx:02d}" + (f"_{safe_name(topic)}" if topic else "")
        lesson_dir = base_dir / f"{prefix}_{safe_name(title)}"
        lesson_dir.mkdir(parents=True, exist_ok=True)
        print(f"  📁 {lesson_dir.name}")

        desc = mat.get("description", "")
        if desc and not (lesson_dir / "description.txt").exists():
            (lesson_dir / "description.txt").write_text(desc, encoding="utf-8")

        process_attachments(
            mat.get("materials", []),
            lesson_dir, slides, drive, creds,
            vimeo_queue, course_title, title,
            extra_text=desc,
        )

    # ── Course work (assignments) ─────────────────────────────────────────────
    print("\n📝 Processing assignments...\n")
    try:
        cws = classroom.courses().courseWork().list(
            courseId=course_id
        ).execute().get("courseWork", [])
    except Exception as e:
        print(f"  ⚠️  Could not fetch coursework: {e}")
        cws = []

    for idx, cw in enumerate(cws, 1):
        title      = cw.get("title", f"Assignment_{idx}")
        topic      = topics.get(cw.get("topicId", ""), "")
        prefix     = f"A{idx:02d}" + (f"_{safe_name(topic)}" if topic else "")
        lesson_dir = base_dir / f"{prefix}_{safe_name(title)}"
        lesson_dir.mkdir(parents=True, exist_ok=True)
        print(f"  📁 {lesson_dir.name}")

        desc = cw.get("description", "")
        if desc and not (lesson_dir / "description.txt").exists():
            (lesson_dir / "description.txt").write_text(desc, encoding="utf-8")

        process_attachments(
            cw.get("materials", []),
            lesson_dir, slides, drive, creds,
            vimeo_queue, course_title, title,
            extra_text=desc,
        )

    # ── Save Vimeo queue (merge with existing, deduplicate by URL) ───────────
    new_videos = [v for v in vimeo_queue if v["url"] not in existing_urls]
    merged_queue = existing_queue + new_videos
    if merged_queue:
        q_path.write_text(
            json.dumps(merged_queue, indent=2, ensure_ascii=False),
            encoding="utf-8"
        )
        if new_videos:
            print(f"\n🎬 {len(new_videos)} new video(s) added → {q_path} ({len(merged_queue)} total)")
        else:
            print(f"\n🎬 No new videos found ({len(merged_queue)} already queued)")
        print("   Run classroom_transcribe.py next.")
    else:
        print("\nℹ️  No Vimeo videos found.")

    print(f"\n🎉 Done! Archived to: {base_dir.resolve()}")


if __name__ == "__main__":
    main()
